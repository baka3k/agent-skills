#!/usr/bin/env python3
"""Generate a full software bidding package from YAML/JSON contracts.

Outputs:
- outputs/proposal/proposal.md
- outputs/estimate/estimate_summary.json
- outputs/staffing/staffing_plan.md
- outputs/quality/quality_gates.md
- outputs/evidence/evidence_log.json
- outputs/slides/bid_deck.pptx (markdown fallback if PPTX runtime unavailable)
"""

from __future__ import annotations

import argparse
import json
import math
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    import yaml
except Exception as exc:  # pragma: no cover
    raise SystemExit(f"PyYAML is required: {exc}")


EXTERNAL_BASELINES = [
    {
        "name": "Gemini image generation docs",
        "url": "https://ai.google.dev/gemini-api/docs/image-generation",
    },
    {
        "name": "Gemini changelog",
        "url": "https://ai.google.dev/gemini-api/docs/changelog",
    },
    {
        "name": "Gemini models catalog",
        "url": "https://ai.google.dev/gemini-api/docs/models",
    },
    {
        "name": "COCOMO II overview",
        "url": "https://boehmcsse.org/tools/cocomo-ii/",
    },
    {
        "name": "ISO/IEC 25010:2023",
        "url": "https://www.iso.org/standard/78176.html",
    },
    {
        "name": "DORA metrics",
        "url": "https://dora.dev/guides/dora-metrics/",
    },
    {
        "name": "NASA Cost Estimating Handbook v4.0",
        "url": "https://www.nasa.gov/ocfo/ppc-corner/nasa-cost-estimating-handbook-ceh/",
    },
]

COMPLEXITY_MULTIPLIER = {
    "low": 0.85,
    "med": 1.0,
    "medium": 1.0,
    "high": 1.25,
}

UNKNOWN_BUFFER = {
    "integration": 0.12,
    "data_migration": 0.10,
    "nfr": 0.08,
}

PHASE_RATIOS = {
    "Discovery": 0.12,
    "Foundation/Architecture": 0.18,
    "Build": 0.50,
    "Stabilize/UAT": 0.15,
    "Go-live/Hypercare": 0.05,
}

PHASE_ROLE_MIX = {
    "Discovery": {
        "BA": 0.25,
        "Presales SA": 0.25,
        "PM": 0.25,
        "Tech Lead": 0.25,
    },
    "Foundation/Architecture": {
        "Architect": 0.30,
        "Lead Dev": 0.30,
        "DevOps": 0.20,
        "QA Lead": 0.20,
    },
    "Build": {
        "Backend Dev": 0.35,
        "Frontend Dev": 0.25,
        "QA Engineer": 0.20,
        "BA": 0.10,
        "PM": 0.10,
    },
    "Stabilize/UAT": {
        "QA Engineer": 0.40,
        "Backend Dev": 0.25,
        "Frontend Dev": 0.15,
        "PM": 0.10,
        "BA": 0.10,
    },
    "Go-live/Hypercare": {
        "SRE/DevOps": 0.35,
        "Lead Dev": 0.25,
        "Support Engineer": 0.25,
        "PM": 0.15,
    },
}


@dataclass
class EstimateSummary:
    base_pm: float
    best_pm: float
    worst_pm: float
    risk_buffer: float
    fixed_price_cost: Optional[Dict[str, Any]]
    tm_cost: Optional[Dict[str, Any]]
    confidence_tier: str


@dataclass
class StaffingSummary:
    duration_months: float
    phase_pm: Dict[str, float]
    role_pm: Dict[str, float]
    monthly_role_fte: Dict[str, float]


def utc_now_iso() -> str:
    return datetime.now(tz=timezone.utc).isoformat()


def load_structured_file(path: Path) -> Any:
    content = path.read_text(encoding="utf-8")
    if path.suffix.lower() == ".json":
        return json.loads(content)
    return yaml.safe_load(content)


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def parse_duration_months(raw_timeline: Any, fallback_pm: float) -> float:
    if not raw_timeline:
        return max(round(fallback_pm / 6.0, 2), 1.0)

    text = str(raw_timeline).lower().strip()
    number_match = re.search(r"(\d+(?:\.\d+)?)", text)
    if not number_match:
        return max(round(fallback_pm / 6.0, 2), 1.0)

    value = float(number_match.group(1))
    if "week" in text:
        return max(round(value / 4.33, 2), 0.5)
    if "day" in text:
        return max(round(value / 30.0, 2), 0.25)
    return max(round(value, 2), 0.5)


def normalize_rate_card(rate_card: Any) -> Dict[str, Any]:
    if rate_card is None:
        return {}

    if isinstance(rate_card, list):
        return {"roles": rate_card}

    if not isinstance(rate_card, dict):
        return {}

    roles = rate_card.get("roles")
    if not isinstance(roles, list):
        roles = []

    return {
        "currency": rate_card.get("currency", "USD"),
        "region": rate_card.get("region", "unknown"),
        "blended_rate_policy": rate_card.get("blended_rate_policy", {}) or {},
        "roles": roles,
    }


def build_rate_lookup(rate_card: Dict[str, Any]) -> Dict[str, float]:
    lookup: Dict[str, float] = {}
    for item in rate_card.get("roles", []):
        role = str(item.get("role", "")).strip()
        monthly_rate = item.get("monthly_rate")
        if not role:
            continue
        try:
            lookup[role] = float(monthly_rate)
        except (TypeError, ValueError):
            continue
    return lookup


def resolve_blended_rate(rate_card: Dict[str, Any], lookup: Dict[str, float]) -> Optional[float]:
    policy = rate_card.get("blended_rate_policy", {}) or {}
    explicit = policy.get("blended_monthly_rate")
    if explicit is not None:
        try:
            return float(explicit)
        except (TypeError, ValueError):
            pass

    if lookup:
        return round(sum(lookup.values()) / len(lookup), 2)

    return None


def normalize_wbs(brief: Dict[str, Any]) -> List[Dict[str, Any]]:
    wbs = brief.get("work_breakdown")
    if isinstance(wbs, list) and wbs:
        normalized: List[Dict[str, Any]] = []
        for item in wbs:
            if not isinstance(item, dict):
                continue
            module = str(item.get("module", "Unnamed module")).strip() or "Unnamed module"
            try:
                base_pm = float(item.get("base_pm", 1.5))
            except (TypeError, ValueError):
                base_pm = 1.5
            complexity = str(item.get("complexity", "medium")).lower()
            if complexity not in COMPLEXITY_MULTIPLIER:
                complexity = "medium"
            normalized.append(
                {
                    "module": module,
                    "base_pm": base_pm,
                    "complexity": complexity,
                    "multiplier": COMPLEXITY_MULTIPLIER[complexity],
                }
            )
        if normalized:
            return normalized

    scope_items = brief.get("scope_in")
    if isinstance(scope_items, list) and scope_items:
        return [
            {
                "module": str(item),
                "base_pm": 1.5,
                "complexity": "medium",
                "multiplier": 1.0,
            }
            for item in scope_items
        ]

    return [
        {
            "module": "Core implementation",
            "base_pm": 5.0,
            "complexity": "medium",
            "multiplier": 1.0,
        }
    ]


def compute_risk_buffer(brief: Dict[str, Any]) -> Tuple[float, List[str]]:
    unknowns = brief.get("unknowns", {})
    assumptions: List[str] = []
    risk_buffer = 0.0

    if isinstance(unknowns, list):
        unknown_map = {str(x): True for x in unknowns}
    elif isinstance(unknowns, dict):
        unknown_map = {str(k): bool(v) for k, v in unknowns.items()}
    else:
        unknown_map = {}

    for key, increment in UNKNOWN_BUFFER.items():
        if unknown_map.get(key, False):
            risk_buffer += increment
            assumptions.append(f"Unknown '{key}' included in risk buffer (+{int(increment * 100)}%).")

    if risk_buffer > 0.30:
        assumptions.append("Risk buffer capped at 30% by policy.")

    return min(risk_buffer, 0.30), assumptions


def compute_confidence(brief: Dict[str, Any], evidence_cfg: Dict[str, Any]) -> Tuple[str, float, List[str]]:
    score = 0.0
    notes: List[str] = []

    has_mind = bool(evidence_cfg.get("mind_source_ids"))
    has_graph = bool(evidence_cfg.get("graph_project_ids"))
    has_web = bool(evidence_cfg.get("internet_domains_allowlist"))
    has_code_paths = bool(brief.get("code_sample_paths"))

    if has_mind:
        score += 0.35
        notes.append("mind_mcp evidence coverage available.")
    if has_graph:
        score += 0.35
        notes.append("graph_mcp evidence coverage available.")
    if has_web:
        score += 0.20
        notes.append("Trusted internet baselines configured.")
    if has_code_paths:
        score += 0.10
        notes.append("Code sample paths provided in brief.")

    unknowns = brief.get("unknowns", {})
    unknown_count = 0
    if isinstance(unknowns, dict):
        unknown_count = sum(1 for value in unknowns.values() if value)
    elif isinstance(unknowns, list):
        unknown_count = len(unknowns)

    if unknown_count > 2:
        score -= 0.10
        notes.append("Confidence reduced due to multiple active unknowns.")

    score = max(min(score, 1.0), 0.0)

    if score >= 0.80:
        return "High", score, notes
    if score >= 0.55:
        return "Medium", score, notes
    return "Low", score, notes


def build_solution_options(brief: Dict[str, Any], base_pm: float, risk_buffer: float) -> List[Dict[str, Any]]:
    scope_in = brief.get("scope_in") or []
    constraints = brief.get("constraints") or []

    baseline = {
        "option": "A - Baseline",
        "focus": "Balanced delivery with lower execution variance",
        "architecture": "Modular service-oriented baseline, conservative integration sequencing",
        "delivery_model": "Milestone-based with stabilization before go-live",
        "impact": {
            "cost": "Baseline",
            "timeline": "Baseline",
            "quality": "Baseline + mandatory quality gates",
            "risk": "Lower implementation volatility",
        },
        "assumptions": [
            f"Scope includes {len(scope_in)} primary streams.",
            f"Risk buffer {round(risk_buffer * 100, 1)}% applied.",
        ],
    }

    optimized = {
        "option": "B - Optimized",
        "focus": "Accelerated delivery via parallelization and platform automation",
        "architecture": "Baseline architecture plus stronger CI/CD and test automation from early phases",
        "delivery_model": "Parallel workstreams with tighter dependency management",
        "impact": {
            "cost": "Slightly higher monthly burn, potentially lower total delay risk",
            "timeline": "Shorter by ~10% when dependencies hold",
            "quality": "Higher test and observability readiness",
            "risk": "Higher coordination demand during build",
        },
        "assumptions": [
            "Client stakeholders provide faster turnaround on integration decisions.",
            "Platform constraints are known before Foundation phase completion.",
            f"Current constraints count: {len(constraints)}.",
        ],
    }

    if base_pm >= 40:
        optimized["assumptions"].append("Program-level governance is required for cross-stream dependencies.")

    return [baseline, optimized]


def compute_estimate(
    brief: Dict[str, Any],
    evidence_cfg: Dict[str, Any],
    rate_card: Dict[str, Any],
) -> Tuple[EstimateSummary, List[Dict[str, Any]], Dict[str, float], List[str], float]:
    assumptions: List[str] = []
    wbs_items = normalize_wbs(brief)

    wbs_pm = 0.0
    for item in wbs_items:
        weighted_pm = float(item["base_pm"]) * float(item["multiplier"])
        item["weighted_pm"] = round(weighted_pm, 3)
        wbs_pm += weighted_pm

    risk_buffer, risk_assumptions = compute_risk_buffer(brief)
    assumptions.extend(risk_assumptions)

    base_pm = round(wbs_pm * (1 + risk_buffer), 2)
    best_pm = round(base_pm * 0.85, 2)
    worst_pm = round(base_pm * 1.30, 2)

    rate_lookup = build_rate_lookup(rate_card)
    blended_rate = resolve_blended_rate(rate_card, rate_lookup)

    fixed_price_cost: Optional[Dict[str, Any]] = None
    if blended_rate is not None:
        fixed_price_cost = {
            "currency": rate_card.get("currency", "USD"),
            "blended_monthly_rate": round(blended_rate, 2),
            "contingency_rate": 0.15,
            "best": round(best_pm * blended_rate * 1.15, 2),
            "base": round(base_pm * blended_rate * 1.15, 2),
            "worst": round(worst_pm * blended_rate * 1.15, 2),
        }
    else:
        assumptions.append("Rate card missing; fixed-price cost is pending.")

    confidence_tier, confidence_score, confidence_notes = compute_confidence(brief, evidence_cfg)
    assumptions.extend(confidence_notes)

    summary = EstimateSummary(
        base_pm=base_pm,
        best_pm=best_pm,
        worst_pm=worst_pm,
        risk_buffer=round(risk_buffer, 4),
        fixed_price_cost=fixed_price_cost,
        tm_cost=None,
        confidence_tier=confidence_tier,
    )

    return summary, wbs_items, rate_lookup, assumptions, confidence_score


def compute_staffing(
    base_pm: float,
    timeline_target: Any,
    rate_lookup: Dict[str, float],
    fallback_rate: Optional[float],
    currency: str,
) -> Tuple[StaffingSummary, Dict[str, Any], List[str]]:
    assumptions: List[str] = []
    duration_months = parse_duration_months(timeline_target, base_pm)

    phase_pm: Dict[str, float] = {}
    role_pm: Dict[str, float] = {}

    for phase, ratio in PHASE_RATIOS.items():
        phase_effort = round(base_pm * ratio, 3)
        phase_pm[phase] = phase_effort

        mix = PHASE_ROLE_MIX[phase]
        for role, pct in mix.items():
            role_pm[role] = round(role_pm.get(role, 0.0) + phase_effort * pct, 3)

    monthly_role_fte: Dict[str, float] = {}
    for role, pm_value in role_pm.items():
        monthly_role_fte[role] = round(pm_value / duration_months, 3)

    tm_role_cost: Dict[str, float] = {}
    tm_base_cost = 0.0
    missing_roles: List[str] = []

    for role, pm_value in role_pm.items():
        role_rate = rate_lookup.get(role)
        if role_rate is None and fallback_rate is not None:
            role_rate = fallback_rate
            missing_roles.append(role)
        if role_rate is None:
            continue
        cost = pm_value * role_rate
        tm_role_cost[role] = round(cost, 2)
        tm_base_cost += cost

    if missing_roles:
        assumptions.append(
            "Fallback blended rate used for roles missing in rate card: " + ", ".join(sorted(missing_roles))
        )

    tm_cost: Dict[str, Any] = {
        "currency": currency,
        "duration_months": round(duration_months, 2),
        "monthly_team_size": round(base_pm / duration_months, 2),
        "role_cost_breakdown": tm_role_cost,
        "best": round(tm_base_cost * 0.85, 2),
        "base": round(tm_base_cost, 2),
        "worst": round(tm_base_cost * 1.30, 2),
    }

    staffing = StaffingSummary(
        duration_months=round(duration_months, 2),
        phase_pm=phase_pm,
        role_pm=role_pm,
        monthly_role_fte=monthly_role_fte,
    )

    return staffing, tm_cost, assumptions


def evaluate_quality_gates(
    brief: Dict[str, Any],
    estimate: EstimateSummary,
    staffing: StaffingSummary,
) -> List[Dict[str, str]]:
    scope_in = brief.get("scope_in") or []
    scope_out = brief.get("scope_out") or []
    goals = brief.get("business_goals") or []
    constraints = brief.get("constraints") or []
    compliance = brief.get("compliance_needs") or []

    gates: List[Dict[str, str]] = []

    scope_pass = bool(scope_in and scope_out and goals)
    gates.append(
        {
            "gate": "Scope clarity",
            "status": "pass" if scope_pass else "warn",
            "evidence": "scope_in/scope_out/business_goals",
            "notes": "Scope boundaries explicit." if scope_pass else "Scope boundaries are incomplete.",
        }
    )

    arch_pass = bool(constraints)
    gates.append(
        {
            "gate": "Architecture readiness",
            "status": "pass" if arch_pass else "warn",
            "evidence": "constraints and solution options",
            "notes": "Core technical constraints provided." if arch_pass else "Constraints are missing or too generic.",
        }
    )

    estimate_pass = estimate.confidence_tier in {"High", "Medium"} and estimate.base_pm > 0
    gates.append(
        {
            "gate": "Estimation confidence",
            "status": "pass" if estimate_pass else "warn",
            "evidence": "estimate range + evidence coverage",
            "notes": f"Confidence tier: {estimate.confidence_tier}.",
        }
    )

    delivery_pass = bool(brief.get("timeline_target")) and staffing.duration_months > 0
    gates.append(
        {
            "gate": "Delivery readiness",
            "status": "pass" if delivery_pass else "warn",
            "evidence": "timeline_target + staffing plan",
            "notes": "Timeline and staffing are aligned." if delivery_pass else "Timeline or staffing alignment missing.",
        }
    )

    prod_pass = bool(compliance)
    gates.append(
        {
            "gate": "Production readiness",
            "status": "pass" if prod_pass else "warn",
            "evidence": "compliance_needs + operations assumptions",
            "notes": "Compliance requirements are captured." if prod_pass else "Compliance requirements not provided.",
        }
    )

    return gates


def build_evidence_log(
    evidence_cfg: Dict[str, Any],
    confidence_tier: str,
    confidence_score: float,
    assumptions: List[str],
) -> List[Dict[str, Any]]:
    now = utc_now_iso()
    entries: List[Dict[str, Any]] = []

    for source_id in evidence_cfg.get("mind_source_ids", []) or []:
        entries.append(
            {
                "source_type": "mind_mcp",
                "source_id": source_id,
                "claim": "Project/domain context retrieved",
                "citation": f"mind_mcp:{source_id}",
                "confidence_weight": 0.35,
                "timestamp": now,
            }
        )

    for project_id in evidence_cfg.get("graph_project_ids", []) or []:
        entries.append(
            {
                "source_type": "graph_mcp",
                "source_id": project_id,
                "claim": "Code structure and dependency context retrieved",
                "citation": f"graph_mcp:project_id={project_id}",
                "confidence_weight": 0.35,
                "timestamp": now,
            }
        )

    allowlist = evidence_cfg.get("internet_domains_allowlist", []) or []
    for baseline in EXTERNAL_BASELINES:
        domain = baseline["url"].split("/")[2]
        if not allowlist or any(domain.endswith(allowed) for allowed in allowlist):
            entries.append(
                {
                    "source_type": "internet",
                    "source_id": domain,
                    "claim": baseline["name"],
                    "citation": baseline["url"],
                    "confidence_weight": 0.20,
                    "timestamp": now,
                }
            )

    entries.append(
        {
            "source_type": "governance",
            "source_id": "confidence_summary",
            "claim": f"Confidence tier={confidence_tier}, score={round(confidence_score, 2)}",
            "citation": "internal-confidence-policy",
            "confidence_weight": round(confidence_score, 3),
            "timestamp": now,
        }
    )

    if assumptions:
        entries.append(
            {
                "source_type": "assumption",
                "source_id": "assumptions",
                "claim": f"{len(assumptions)} assumptions were recorded",
                "citation": "proposal-assumptions",
                "confidence_weight": 0.0,
                "timestamp": now,
            }
        )

    return entries


def to_md_bullets(items: List[Any]) -> str:
    if not items:
        return "- None"
    return "\n".join(f"- {str(item)}" for item in items)


def render_solution_md(solution_options: List[Dict[str, Any]]) -> str:
    lines: List[str] = ["# Solution Options", ""]
    for option in solution_options:
        lines.append(f"## {option['option']}")
        lines.append(f"Focus: {option['focus']}")
        lines.append(f"Architecture: {option['architecture']}")
        lines.append(f"Delivery model: {option['delivery_model']}")
        lines.append("")
        lines.append("Tradeoff impact:")
        for key, value in option["impact"].items():
            lines.append(f"- {key}: {value}")
        lines.append("")
        lines.append("Assumptions:")
        for assumption in option["assumptions"]:
            lines.append(f"- {assumption}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def render_staffing_md(staffing: StaffingSummary) -> str:
    phase_lines = [
        "| Phase | Effort (PM) |",
        "| --- | ---: |",
    ]
    for phase, value in staffing.phase_pm.items():
        phase_lines.append(f"| {phase} | {value:.2f} |")

    role_lines = [
        "| Role | PM | Avg Monthly FTE |",
        "| --- | ---: | ---: |",
    ]
    for role in sorted(staffing.role_pm.keys()):
        role_lines.append(
            f"| {role} | {staffing.role_pm[role]:.2f} | {staffing.monthly_role_fte[role]:.2f} |"
        )

    return "\n".join(
        [
            "# Staffing Plan",
            "",
            f"Estimated duration: **{staffing.duration_months:.2f} months**",
            "",
            "## Phase Allocation",
            "",
            *phase_lines,
            "",
            "## Role Allocation",
            "",
            *role_lines,
            "",
        ]
    )


def render_quality_md(gates: List[Dict[str, str]]) -> str:
    lines = [
        "# Quality Gates",
        "",
        "| Gate | Status | Evidence | Notes |",
        "| --- | --- | --- | --- |",
    ]
    for gate in gates:
        lines.append(
            f"| {gate['gate']} | {gate['status']} | {gate['evidence']} | {gate['notes']} |"
        )
    lines.append("")
    return "\n".join(lines)


def render_proposal_md(
    brief: Dict[str, Any],
    estimate: EstimateSummary,
    staffing: StaffingSummary,
    solution_options: List[Dict[str, Any]],
    gates: List[Dict[str, str]],
    assumptions: List[str],
    evidence_log: List[Dict[str, Any]],
) -> str:
    client_name = (brief.get("client_context") or {}).get("client_name", "Client")

    fixed_price = (
        "Cost pending (rate card not provided)."
        if estimate.fixed_price_cost is None
        else (
            f"{estimate.fixed_price_cost['currency']} "
            f"best/base/worst = {estimate.fixed_price_cost['best']:.2f} / "
            f"{estimate.fixed_price_cost['base']:.2f} / {estimate.fixed_price_cost['worst']:.2f}"
        )
    )

    tm_cost = (
        "Cost pending (rate card not provided)."
        if estimate.tm_cost is None
        else (
            f"{estimate.tm_cost['currency']} "
            f"best/base/worst = {estimate.tm_cost['best']:.2f} / "
            f"{estimate.tm_cost['base']:.2f} / {estimate.tm_cost['worst']:.2f}"
        )
    )

    gate_summary = ", ".join(f"{g['gate']}: {g['status']}" for g in gates)

    citations = sorted(
        {
            entry["citation"]
            for entry in evidence_log
            if entry.get("source_type") in {"internet", "mind_mcp", "graph_mcp"}
        }
    )

    return "\n".join(
        [
            f"# Project Proposal - {client_name}",
            "",
            "## Executive Summary",
            "",
            "This proposal provides a hybrid bidding package with fixed-price and T&M views, "
            "including estimation ranges, staffing strategy, quality gates, and traceable evidence.",
            "",
            "## Scope",
            "",
            "### In Scope",
            to_md_bullets(brief.get("scope_in") or []),
            "",
            "### Out of Scope",
            to_md_bullets(brief.get("scope_out") or []),
            "",
            "## Solution Options",
            "",
            f"- {solution_options[0]['option']}: {solution_options[0]['focus']}",
            f"- {solution_options[1]['option']}: {solution_options[1]['focus']}",
            "",
            "## Estimation Summary",
            "",
            f"- Effort (PM) best/base/worst: {estimate.best_pm:.2f} / {estimate.base_pm:.2f} / {estimate.worst_pm:.2f}",
            f"- Risk buffer: {estimate.risk_buffer * 100:.1f}%",
            f"- Confidence tier: {estimate.confidence_tier}",
            f"- Fixed-price: {fixed_price}",
            f"- T&M: {tm_cost}",
            "",
            "## Staffing Strategy",
            "",
            f"- Delivery duration assumption: {staffing.duration_months:.2f} months",
            "- Phases: Discovery, Foundation/Architecture, Build, Stabilize/UAT, Go-live/Hypercare",
            "",
            "## Quality Gates",
            "",
            f"- {gate_summary}",
            "",
            "## Assumptions",
            "",
            to_md_bullets(assumptions),
            "",
            "## Citations",
            "",
            to_md_bullets(citations),
            "",
        ]
    )


def try_generate_pptx(
    output_file: Path,
    brief: Dict[str, Any],
    estimate: EstimateSummary,
    staffing: StaffingSummary,
    gates: List[Dict[str, str]],
    evidence_log: List[Dict[str, Any]],
    solution_options: List[Dict[str, Any]],
) -> Tuple[bool, str]:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        return False, f"python-pptx not available: {exc}"

    prs = Presentation()

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets_slide(title: str, bullets: List[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        if not bullets:
            bullets = ["No content"]
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet

    client = (brief.get("client_context") or {}).get("client_name", "Client")
    add_title_slide(
        f"Software Project Bid - {client}",
        f"Generated {datetime.now().strftime('%Y-%m-%d')} | Hybrid (Fixed-price + T&M)",
    )

    add_bullets_slide(
        "Scope Overview",
        [
            "In scope:",
            *[f"- {x}" for x in (brief.get("scope_in") or [])],
            "Out of scope:",
            *[f"- {x}" for x in (brief.get("scope_out") or [])],
        ],
    )

    add_bullets_slide(
        "Solution Options",
        [
            f"{solution_options[0]['option']}: {solution_options[0]['focus']}",
            f"{solution_options[1]['option']}: {solution_options[1]['focus']}",
            "Tradeoffs cover cost, timeline, quality, and delivery risk.",
        ],
    )

    fixed_text = (
        "Fixed-price cost pending (no rate card)."
        if not estimate.fixed_price_cost
        else (
            f"Fixed-price best/base/worst: {estimate.fixed_price_cost['best']:.2f}/"
            f"{estimate.fixed_price_cost['base']:.2f}/{estimate.fixed_price_cost['worst']:.2f} "
            f"{estimate.fixed_price_cost['currency']}"
        )
    )

    tm_text = (
        "T&M cost pending (no rate card)."
        if not estimate.tm_cost
        else (
            f"T&M best/base/worst: {estimate.tm_cost['best']:.2f}/"
            f"{estimate.tm_cost['base']:.2f}/{estimate.tm_cost['worst']:.2f} "
            f"{estimate.tm_cost['currency']}"
        )
    )

    add_bullets_slide(
        "Estimate and Cost Range",
        [
            f"Effort PM best/base/worst: {estimate.best_pm:.2f}/{estimate.base_pm:.2f}/{estimate.worst_pm:.2f}",
            f"Risk buffer: {estimate.risk_buffer * 100:.1f}%",
            f"Confidence: {estimate.confidence_tier}",
            fixed_text,
            tm_text,
        ],
    )

    phase_lines = [f"{phase}: {pm:.2f} PM" for phase, pm in staffing.phase_pm.items()]
    add_bullets_slide("Staffing by Phase", phase_lines)

    gate_lines = [f"{g['gate']}: {g['status']} ({g['notes']})" for g in gates]
    add_bullets_slide("Quality Gates", gate_lines)

    citations = sorted(
        {
            x["citation"]
            for x in evidence_log
            if x.get("source_type") in {"internet", "mind_mcp", "graph_mcp"}
        }
    )
    add_bullets_slide("Citations", citations[:12] or ["No citations"])

    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    return True, "PPTX generated"


def render_slide_markdown_fallback(output_file: Path, message: str) -> None:
    output_file.parent.mkdir(parents=True, exist_ok=True)
    output_file.write_text(
        "\n".join(
            [
                "# Bid Deck Fallback",
                "",
                "PPTX generation was not available in this run.",
                "",
                f"Reason: {message}",
                "",
                "Use $bid-slide-factory with $presentations to generate the final deck.",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def write_json(path: Path, data: Any) -> None:
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Generate a full software bidding package.")
    parser.add_argument("--bid-brief", required=True, help="Path to bid_brief.yaml/json")
    parser.add_argument("--evidence-config", required=True, help="Path to evidence_config.yaml/json")
    parser.add_argument("--rate-card", default=None, help="Path to rate_card.yaml/json")
    parser.add_argument("--output-dir", default="outputs", help="Output directory")
    parser.add_argument("--skip-slides", action="store_true", help="Skip PPTX generation")
    return parser


def main() -> int:
    args = build_parser().parse_args()

    bid_brief_path = Path(args.bid_brief).expanduser().resolve()
    evidence_cfg_path = Path(args.evidence_config).expanduser().resolve()
    rate_card_path = Path(args.rate_card).expanduser().resolve() if args.rate_card else None
    output_dir = Path(args.output_dir).expanduser().resolve()

    brief = load_structured_file(bid_brief_path) or {}
    evidence_cfg = load_structured_file(evidence_cfg_path) or {}
    rate_card_raw = load_structured_file(rate_card_path) if rate_card_path else None
    rate_card = normalize_rate_card(rate_card_raw)

    estimate, wbs_items, rate_lookup, assumptions, confidence_score = compute_estimate(
        brief, evidence_cfg, rate_card
    )

    blended_rate = resolve_blended_rate(rate_card, rate_lookup)
    staffing, tm_cost, staffing_assumptions = compute_staffing(
        estimate.base_pm,
        brief.get("timeline_target"),
        rate_lookup,
        blended_rate,
        str(rate_card.get("currency", "USD")),
    )
    assumptions.extend(staffing_assumptions)

    estimate.tm_cost = tm_cost if rate_lookup or blended_rate else None
    if estimate.tm_cost is None:
        assumptions.append("Rate card missing; T&M cost is pending.")

    solution_options = build_solution_options(brief, estimate.base_pm, estimate.risk_buffer)
    gates = evaluate_quality_gates(brief, estimate, staffing)
    evidence_log = build_evidence_log(
        evidence_cfg,
        estimate.confidence_tier,
        confidence_score,
        assumptions,
    )

    # Output paths
    proposal_dir = output_dir / "proposal"
    estimate_dir = output_dir / "estimate"
    staffing_dir = output_dir / "staffing"
    quality_dir = output_dir / "quality"
    evidence_dir = output_dir / "evidence"
    slides_dir = output_dir / "slides"
    solution_dir = output_dir / "solution"

    for p in [proposal_dir, estimate_dir, staffing_dir, quality_dir, evidence_dir, slides_dir, solution_dir]:
        ensure_dir(p)

    proposal_md = render_proposal_md(
        brief,
        estimate,
        staffing,
        solution_options,
        gates,
        assumptions,
        evidence_log,
    )
    (proposal_dir / "proposal.md").write_text(proposal_md, encoding="utf-8")

    estimate_payload = {
        "generated_at": utc_now_iso(),
        "delivery_mode": brief.get("delivery_mode", "hybrid"),
        "target_language": brief.get("target_language", "en"),
        "wbs": wbs_items,
        "effort_pm": {
            "best": estimate.best_pm,
            "base": estimate.base_pm,
            "worst": estimate.worst_pm,
            "risk_buffer": estimate.risk_buffer,
        },
        "cost": {
            "fixed_price": estimate.fixed_price_cost,
            "tm": estimate.tm_cost,
        },
        "confidence": {
            "tier": estimate.confidence_tier,
            "score": round(confidence_score, 3),
        },
        "assumptions": assumptions,
    }
    write_json(estimate_dir / "estimate_summary.json", estimate_payload)

    (staffing_dir / "staffing_plan.md").write_text(render_staffing_md(staffing), encoding="utf-8")
    (quality_dir / "quality_gates.md").write_text(render_quality_md(gates), encoding="utf-8")
    write_json(evidence_dir / "evidence_log.json", evidence_log)
    (solution_dir / "solution_options.md").write_text(render_solution_md(solution_options), encoding="utf-8")

    slides_note = "Slide generation skipped by flag"
    if args.skip_slides:
        render_slide_markdown_fallback(slides_dir / "bid_deck.md", slides_note)
    else:
        success, message = try_generate_pptx(
            slides_dir / "bid_deck.pptx",
            brief,
            estimate,
            staffing,
            gates,
            evidence_log,
            solution_options,
        )
        if not success:
            render_slide_markdown_fallback(slides_dir / "bid_deck.md", message)

    summary = {
        "status": "ok",
        "output_dir": str(output_dir),
        "proposal": str(proposal_dir / "proposal.md"),
        "estimate": str(estimate_dir / "estimate_summary.json"),
        "staffing": str(staffing_dir / "staffing_plan.md"),
        "quality": str(quality_dir / "quality_gates.md"),
        "evidence": str(evidence_dir / "evidence_log.json"),
        "slides": str(slides_dir / "bid_deck.pptx"),
    }
    print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
